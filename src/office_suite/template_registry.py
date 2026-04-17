"""File-based template registry for slides and documents.

This module exposes a simple, agent-friendly catalog of HTML templates stored
on disk (`src/office_suite/templates/`). Each template is a folder with:

- `template.html` — Jinja2 template
- `meta.json` — declarative metadata with the slots the agent must fill

Agents (Hermes / OpenClaw) can:

1. Call `list_templates()` to discover available templates and their slots.
2. Call `get_template_meta(template_id)` for a single template contract.
3. Call `render_slide_template(...)` / `render_document_template(...)` to
   produce `.html`, `.png`, `.pdf`, or `.pptx` deliverables.

The registry is intentionally data-driven: to add a new template, drop a new
folder in `templates/slides/<slug>/` or `templates/documents/<slug>/`. No
Python changes required.
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from jinja2 import ChainableUndefined, Environment, FileSystemLoader, select_autoescape

TEMPLATES_ROOT = Path(__file__).resolve().parent / "templates"
SLIDES_ROOT = TEMPLATES_ROOT / "slides"
DOCUMENTS_ROOT = TEMPLATES_ROOT / "documents"

SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080


def _jinja_env(search_paths: Iterable[Path]) -> Environment:
    env = Environment(
        loader=FileSystemLoader([str(p) for p in search_paths]),
        autoescape=select_autoescape(["html", "xml"]),
        undefined=ChainableUndefined,
        trim_blocks=True,
        lstrip_blocks=True,
    )
    return env


def _load_meta(template_dir: Path) -> Dict[str, Any]:
    meta_file = template_dir / "meta.json"
    if not meta_file.exists():
        raise FileNotFoundError(f"meta.json missing for template: {template_dir}")
    meta = json.loads(meta_file.read_text(encoding="utf-8"))
    meta.setdefault("id", template_dir.name)
    meta.setdefault("dir", str(template_dir))
    meta.setdefault("category", template_dir.parent.name)
    return meta


def _iter_templates(root: Path) -> List[Path]:
    if not root.exists():
        return []
    return sorted(p for p in root.iterdir() if p.is_dir() and (p / "template.html").exists())


def list_templates(category: Optional[str] = None) -> Dict[str, Any]:
    """Return a compact catalog of available templates.

    Args:
        category: `"slides"` | `"documents"` | None (both).

    Returns:
        {"slides": [meta, ...], "documents": [meta, ...]}
    """
    out: Dict[str, Any] = {"slides": [], "documents": []}
    if category in (None, "slides"):
        out["slides"] = [_load_meta(p) for p in _iter_templates(SLIDES_ROOT)]
    if category in (None, "documents"):
        out["documents"] = [_load_meta(p) for p in _iter_templates(DOCUMENTS_ROOT)]
    return out


def _resolve_template_dir(template_id: str, category: str) -> Path:
    root = SLIDES_ROOT if category == "slides" else DOCUMENTS_ROOT
    candidate = root / template_id
    if not candidate.exists():
        raise FileNotFoundError(
            f"Template '{template_id}' not found in {category}. "
            f"Known: {[p.name for p in _iter_templates(root)]}"
        )
    return candidate


def get_template_meta(template_id: str, category: str = "slides") -> Dict[str, Any]:
    """Return the `meta.json` contract for a single template."""
    return _load_meta(_resolve_template_dir(template_id, category))


def _merge_defaults(meta: Dict[str, Any], data: Dict[str, Any]) -> Dict[str, Any]:
    defaults: Dict[str, Any] = {}
    for slot in meta.get("slots", []):
        if "default" in slot:
            defaults[slot["name"]] = slot["default"]
    merged = {**defaults, **(data or {})}
    return merged


def render_template_to_html(
    template_id: str,
    data: Dict[str, Any],
    category: str = "slides",
) -> str:
    """Render a single template to an HTML string."""
    template_dir = _resolve_template_dir(template_id, category)
    meta = _load_meta(template_dir)
    env = _jinja_env([template_dir, TEMPLATES_ROOT / "_partials"])
    template = env.get_template("template.html")
    merged = _merge_defaults(meta, data)
    return template.render(**merged, _meta=meta)


def _chrome_bin() -> str:
    for candidate in ("google-chrome", "google-chrome-stable", "chromium", "chromium-browser"):
        if shutil.which(candidate):
            return candidate
    raise RuntimeError(
        "No headless Chrome binary found. Install google-chrome or chromium, "
        "or set OFFICE_CHROME_BIN to the binary path."
    )


def _html_to_png(html_file: Path, out_png: Path, width: int, height: int) -> None:
    out_png.parent.mkdir(parents=True, exist_ok=True)
    chrome = os.environ.get("OFFICE_CHROME_BIN") or _chrome_bin()
    cmd = [
        chrome,
        "--headless",
        "--disable-gpu",
        "--no-sandbox",
        "--hide-scrollbars",
        "--force-device-scale-factor=1",
        f"--screenshot={str(out_png)}",
        f"--window-size={width},{height}",
        html_file.resolve().as_uri(),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)


def _html_to_pdf(html_file: Path, out_pdf: Path) -> None:
    out_pdf.parent.mkdir(parents=True, exist_ok=True)
    chrome = os.environ.get("OFFICE_CHROME_BIN") or _chrome_bin()
    cmd = [
        chrome,
        "--headless",
        "--disable-gpu",
        "--no-sandbox",
        "--hide-scrollbars",
        f"--print-to-pdf={str(out_pdf)}",
        "--no-pdf-header-footer",
        "--print-to-pdf-no-header",
        html_file.resolve().as_uri(),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)


def render_slide_template(
    template_id: str,
    data: Dict[str, Any],
    output_path: str,
    output_format: str = "png",
    width: int = SLIDE_WIDTH_PX,
    height: int = SLIDE_HEIGHT_PX,
) -> Dict[str, Any]:
    """Render ONE slide template to `png` or `html`.

    Most agents use this to preview a single slide. To build a full PPTX,
    prefer `build_pptx_from_templates()`.
    """
    html = render_template_to_html(template_id, data, category="slides")
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    if output_format == "html":
        Path(output_path).write_text(html, encoding="utf-8")
        return {"output_path": output_path, "format": "html", "template": template_id}

    with tempfile.TemporaryDirectory(prefix="docflow_tpl_") as td:
        html_file = Path(td) / "slide.html"
        html_file.write_text(html, encoding="utf-8")
        if output_format == "png":
            _html_to_png(html_file, Path(output_path), width, height)
            return {"output_path": output_path, "format": "png", "template": template_id}
        if output_format == "pdf":
            _html_to_pdf(html_file, Path(output_path))
            return {"output_path": output_path, "format": "pdf", "template": template_id}

    raise ValueError(f"Unsupported output_format: {output_format}")


def build_pptx_from_templates(
    slides: List[Dict[str, Any]],
    output_path: str,
    keep_renders: bool = False,
    renders_dir: Optional[str] = None,
    width: int = SLIDE_WIDTH_PX,
    height: int = SLIDE_HEIGHT_PX,
) -> Dict[str, Any]:
    """Build a PPTX file from a list of templated slides.

    Each element of `slides` must be:
        {"template": "<template_id>", "data": {...}}
    """
    from pptx import Presentation
    from pptx.util import Inches

    if not slides:
        raise ValueError("slides must be a non-empty list")

    with tempfile.TemporaryDirectory(prefix="docflow_tpl_deck_") as td:
        temp_dir = Path(td)
        png_paths: List[Path] = []
        html_paths: List[Path] = []

        for i, item in enumerate(slides, start=1):
            tpl_id = item.get("template")
            data = item.get("data", {}) or {}
            if not tpl_id:
                raise ValueError(f"slides[{i - 1}].template is required")

            html = render_template_to_html(tpl_id, data, category="slides")
            html_file = temp_dir / f"slide_{i:02d}.html"
            png_file = temp_dir / f"slide_{i:02d}.png"
            html_file.write_text(html, encoding="utf-8")
            _html_to_png(html_file, png_file, width, height)
            html_paths.append(html_file)
            png_paths.append(png_file)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for png in png_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                str(png), Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height,
            )

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)

        kept: List[str] = []
        if keep_renders:
            target = Path(renders_dir) if renders_dir else Path(output_path).with_suffix("")
            target.mkdir(parents=True, exist_ok=True)
            for hp, pp in zip(html_paths, png_paths):
                (target / hp.name).write_text(hp.read_text(encoding="utf-8"), encoding="utf-8")
                (target / pp.name).write_bytes(pp.read_bytes())
                kept.extend([str(target / hp.name), str(target / pp.name)])

    return {
        "output_path": output_path,
        "file_size": os.path.getsize(output_path),
        "slides_count": len(slides),
        "render_engine": "html-css-headless-chrome",
        "kept_render_files": kept,
    }


def render_document_template(
    template_id: str,
    data: Dict[str, Any],
    output_path: str,
    output_format: str = "pdf",
) -> Dict[str, Any]:
    """Render a document template (invoice, report, letter, etc.) to pdf/html/png.

    PDFs are produced directly by headless Chrome, so they preserve CSS faithfully.
    """
    html = render_template_to_html(template_id, data, category="documents")
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    if output_format == "html":
        Path(output_path).write_text(html, encoding="utf-8")
        return {"output_path": output_path, "format": "html", "template": template_id}

    with tempfile.TemporaryDirectory(prefix="docflow_doc_") as td:
        html_file = Path(td) / "doc.html"
        html_file.write_text(html, encoding="utf-8")
        if output_format == "pdf":
            _html_to_pdf(html_file, Path(output_path))
            return {"output_path": output_path, "format": "pdf", "template": template_id}
        if output_format == "png":
            _html_to_png(html_file, Path(output_path), 1240, 1754)
            return {"output_path": output_path, "format": "png", "template": template_id}

    raise ValueError(f"Unsupported output_format: {output_format}")


__all__ = [
    "list_templates",
    "get_template_meta",
    "render_template_to_html",
    "render_slide_template",
    "build_pptx_from_templates",
    "render_document_template",
    "TEMPLATES_ROOT",
]
