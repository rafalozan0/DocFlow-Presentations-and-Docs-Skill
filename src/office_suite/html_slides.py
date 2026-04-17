import json
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List

from jinja2 import Template
from pptx import Presentation
from pptx.util import Inches


REFERENCE_THEMES: Dict[str, Dict[str, Any]] = {
    "neon-executive": {
        "background": "#07090F",
        "surface": "#111827",
        "text": "#F5F7FF",
        "muted": "#A9B2C7",
        "accent": "#4D6BFF",
        "accent2": "#7A4DFF",
        "font": "Inter, Manrope, 'Segoe UI', Arial, sans-serif",
    },
    "cobalt-corporate": {
        "background": "#F7F9FC",
        "surface": "#FFFFFF",
        "text": "#0F172A",
        "muted": "#475569",
        "accent": "#2F66D8",
        "accent2": "#1E40AF",
        "font": "Inter, 'Segoe UI', Arial, sans-serif",
    },
    "bold-pitch": {
        "background": "#000000",
        "surface": "#111111",
        "text": "#F8FAFC",
        "muted": "#9CA3AF",
        "accent": "#00D1B2",
        "accent2": "#FFE500",
        "font": "Sora, Poppins, 'Segoe UI', Arial, sans-serif",
    },
}


SLIDE_TEMPLATE = Template(
    """
<!doctype html>
<html>
<head>
  <meta charset=\"utf-8\" />
  <style>
    * { box-sizing: border-box; }
    html, body {
      margin: 0;
      width: 1920px;
      height: 1080px;
      background: {{ theme.background }};
      color: {{ theme.text }};
      font-family: {{ theme.font }};
    }
    .slide {
      width: 1920px;
      height: 1080px;
      position: relative;
      overflow: hidden;
      padding: 72px 96px;
      background: {{ theme.background }};
    }
    .bar {
      position: absolute;
      left: 0;
      top: 0;
      width: 100%;
      height: 14px;
      background: linear-gradient(90deg, {{ theme.accent }}, {{ theme.accent2 }});
    }
    .kicker {
      font-size: 24px;
      letter-spacing: 0.08em;
      text-transform: uppercase;
      color: {{ theme.muted }};
      margin-bottom: 20px;
    }
    h1 {
      margin: 0;
      font-size: 88px;
      line-height: 0.95;
      letter-spacing: -0.02em;
      max-width: 1360px;
    }
    .subtitle {
      margin-top: 30px;
      max-width: 1200px;
      font-size: 34px;
      line-height: 1.25;
      color: {{ theme.muted }};
    }
    .layout-2col {
      display: grid;
      grid-template-columns: 1.05fr 0.95fr;
      gap: 40px;
      margin-top: 36px;
      height: 820px;
    }
    .panel {
      background: {{ theme.surface }};
      border: 1px solid color-mix(in srgb, {{ theme.text }} 14%, transparent);
      border-radius: 20px;
      padding: 28px 30px;
      overflow: hidden;
    }
    .panel h2 {
      margin: 0 0 14px 0;
      font-size: 44px;
      line-height: 1.05;
    }
    .panel p, .panel li {
      font-size: 29px;
      line-height: 1.35;
      color: {{ theme.muted }};
      margin: 0 0 12px 0;
    }
    ul { margin: 0; padding-left: 28px; }
    .stat {
      margin-top: 14px;
      font-size: 96px;
      line-height: 1;
      font-weight: 700;
      color: {{ theme.accent }};
    }
    .footer {
      position: absolute;
      left: 96px;
      right: 96px;
      bottom: 40px;
      display: flex;
      justify-content: space-between;
      font-size: 22px;
      color: {{ theme.muted }};
    }
    .chip {
      border: 1px solid color-mix(in srgb, {{ theme.text }} 22%, transparent);
      border-radius: 999px;
      padding: 8px 16px;
      display: inline-block;
      margin-right: 10px;
      margin-top: 12px;
      font-size: 20px;
      color: {{ theme.text }};
    }
    .shape {
      position: absolute;
      border-radius: 999px;
      filter: blur(2px);
      opacity: 0.36;
      pointer-events: none;
    }
    .shape.a {
      width: 420px;
      height: 420px;
      background: {{ theme.accent }};
      right: -120px;
      top: 120px;
    }
    .shape.b {
      width: 280px;
      height: 280px;
      background: {{ theme.accent2 }};
      right: 300px;
      bottom: -80px;
    }
  </style>
</head>
<body>
  <section class=\"slide\">
    <div class=\"bar\"></div>
    <div class=\"shape a\"></div>
    <div class=\"shape b\"></div>

    <div class=\"kicker\">{{ slide.kicker }}</div>
    <h1>{{ slide.title }}</h1>
    {% if slide.subtitle %}
      <div class=\"subtitle\">{{ slide.subtitle }}</div>
    {% endif %}

    {% if slide.layout == 'two-column' %}
    <div class=\"layout-2col\">
      <div class=\"panel\">
        <h2>{{ slide.left_title }}</h2>
        {% if slide.left_list %}
        <ul>
          {% for item in slide.left_list %}
          <li>{{ item }}</li>
          {% endfor %}
        </ul>
        {% else %}
        <p>{{ slide.left_text }}</p>
        {% endif %}
      </div>
      <div class=\"panel\">
        <h2>{{ slide.right_title }}</h2>
        <p>{{ slide.right_text }}</p>
        {% if slide.stat %}
          <div class=\"stat\">{{ slide.stat }}</div>
        {% endif %}
        {% if slide.tags %}
          <div>
            {% for t in slide.tags %}<span class=\"chip\">{{ t }}</span>{% endfor %}
          </div>
        {% endif %}
      </div>
    </div>
    {% endif %}

    <div class=\"footer\">
      <span>{{ slide.footer_left }}</span>
      <span>{{ slide.footer_right }}</span>
    </div>
  </section>
</body>
</html>
"""
)


def list_reference_themes() -> List[str]:
    return sorted(REFERENCE_THEMES.keys())


def _render_slide_html(slide: Dict[str, Any], theme: Dict[str, Any], out_html: Path) -> None:
    html = SLIDE_TEMPLATE.render(slide=slide, theme=theme)
    out_html.write_text(html, encoding="utf-8")


def _html_to_png(html_file: Path, out_png: Path) -> None:
    out_png.parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        "google-chrome",
        "--headless",
        "--disable-gpu",
        "--no-sandbox",
        f"--screenshot={str(out_png)}",
        "--window-size=1920,1080",
        html_file.resolve().as_uri(),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)


def build_pptx_from_html_theme(
    title: str,
    slides: List[Dict[str, Any]],
    output_path: str,
    theme_name: str = "neon-executive",
    keep_renders: bool = False,
    renders_dir: str = "",
) -> Dict[str, Any]:
    theme = REFERENCE_THEMES.get(theme_name)
    if not theme:
        raise ValueError(f"Unsupported html theme: {theme_name}")

    with tempfile.TemporaryDirectory(prefix="docflow_html_slides_") as td:
        temp_dir = Path(td)
        png_paths: List[Path] = []
        html_paths: List[Path] = []

        for i, slide in enumerate(slides, start=1):
            slide_payload = {
                "kicker": slide.get("kicker", f"Slide {i:02d}"),
                "title": slide.get("title", title),
                "subtitle": slide.get("subtitle", ""),
                "layout": slide.get("layout", "two-column"),
                "left_title": slide.get("left_title", "Context"),
                "left_text": slide.get("left_text", ""),
                "left_list": slide.get("left_list", []),
                "right_title": slide.get("right_title", "Outcome"),
                "right_text": slide.get("right_text", ""),
                "stat": slide.get("stat", ""),
                "tags": slide.get("tags", []),
                "footer_left": slide.get("footer_left", "DocFlow"),
                "footer_right": slide.get("footer_right", f"{i}/{len(slides)}"),
            }

            html_file = temp_dir / f"slide_{i:02d}.html"
            png_file = temp_dir / f"slide_{i:02d}.png"
            _render_slide_html(slide_payload, theme, html_file)
            _html_to_png(html_file, png_file)
            html_paths.append(html_file)
            png_paths.append(png_file)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for png in png_paths:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            s.shapes.add_picture(str(png), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)

        kept_files: List[str] = []
        if keep_renders:
            target_dir = Path(renders_dir) if renders_dir else Path(output_path).with_suffix("")
            target_dir.mkdir(parents=True, exist_ok=True)
            for hp, pp in zip(html_paths, png_paths):
                h_to = target_dir / hp.name
                p_to = target_dir / pp.name
                h_to.write_text(hp.read_text(encoding="utf-8"), encoding="utf-8")
                p_to.write_bytes(pp.read_bytes())
                kept_files.extend([str(h_to), str(p_to)])

    return {
        "output_path": output_path,
        "file_size": os.path.getsize(output_path),
        "slides_count": len(slides),
        "render_engine": "html-css-google-chrome-headless",
        "html_theme": theme_name,
        "kept_render_files": kept_files,
    }


def load_slides_json(path: str) -> List[Dict[str, Any]]:
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    if not isinstance(data, list):
        raise ValueError("Slide JSON must be a list of slide objects")
    return data
