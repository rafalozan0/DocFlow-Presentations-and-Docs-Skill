import os
import tempfile
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

THEMES: Dict[str, Dict[str, Any]] = {
    "midnight-luxe": {
        "background": "0B1020",
        "title_color": "F8FAFC",
        "body_color": "CBD5E1",
        "accent": "7C3AED",
        "muted": "94A3B8",
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "aurora-glow": {
        "background": "F8FAFF",
        "title_color": "0F172A",
        "body_color": "334155",
        "accent": "06B6D4",
        "muted": "64748B",
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "obsidian-slate": {
        "background": "111827",
        "title_color": "F9FAFB",
        "body_color": "D1D5DB",
        "accent": "22D3EE",
        "muted": "9CA3AF",
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "ivory-bloom": {
        "background": "FFFDF7",
        "title_color": "1F2937",
        "body_color": "475569",
        "accent": "F59E0B",
        "muted": "94A3B8",
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "neon-velocity": {
        "background": "0A0F1F",
        "title_color": "E2E8F0",
        "body_color": "CBD5E1",
        "accent": "22C55E",
        "muted": "94A3B8",
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
}

TONE_CONFIG: Dict[str, Dict[str, Any]] = {
    "classic-formal": {"title_size": 36, "body_size": 20, "bullet": "•"},
    "boardroom": {"title_size": 40, "body_size": 20, "bullet": "▸"},
    "conversational": {"title_size": 36, "body_size": 20, "bullet": "•"},
    "laid-back": {"title_size": 34, "body_size": 19, "bullet": "–"},
}


def list_themes() -> List[str]:
    return sorted(THEMES.keys())


def list_chart_modes() -> List[str]:
    return ["native", "matplotlib", "auto"]


def list_tones() -> List[str]:
    return sorted(TONE_CONFIG.keys())


def presentation_preflight_wizard(input_fn=input, print_fn=print) -> Dict[str, Any]:
    """Interactive preflight to collect deck preferences before generation."""

    def choose(label: str, options: List[str], default: str) -> str:
        print_fn(f"\n{label}")
        for i, opt in enumerate(options, 1):
            print_fn(f"  {i}. {opt}")
        raw = input_fn(f"Choose [{default}]: ").strip()
        if not raw:
            return default
        if raw.isdigit():
            idx = int(raw) - 1
            if 0 <= idx < len(options):
                return options[idx]
            return default
        return raw if raw in options else default

    theme = choose("Theme", list_themes(), "midnight-luxe")
    chart_mode = choose("Chart mode", list_chart_modes(), "native")
    tone = choose("Language tone", list_tones(), "boardroom")

    use_emojis_raw = input_fn("Use emojis? [y/N]: ").strip().lower()
    use_emojis = use_emojis_raw in {"y", "yes", "1", "true"}

    return {
        "theme": theme,
        "chart_mode": chart_mode,
        "tone": tone,
        "use_emojis": use_emojis,
    }


def _rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip().lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _apply_background(slide, theme: Dict[str, Any]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme["background"])


def _add_accent_bar(slide, theme: Dict[str, Any], width_inches: float = 13.333) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(width_inches),
        Inches(0.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme["accent"])
    shape.line.fill.background()


def _style_title(slide, layout_type: str, theme: Dict[str, Any], tone_cfg: Dict[str, Any]) -> None:
    if not slide.shapes.title:
        return
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.name = theme["title_font"]
        para.font.size = Pt(44 if layout_type == "title" else tone_cfg["title_size"])
        para.font.bold = True
        para.font.color.rgb = _rgb(theme["title_color"])
        para.alignment = PP_ALIGN.CENTER if layout_type == "title" else PP_ALIGN.LEFT


def _emoji_for_index(idx: int) -> str:
    emojis = ["🚀", "📈", "✅", "💡", "⚙️", "🎯", "📊", "🧠"]
    return emojis[idx % len(emojis)]


def _format_line(line: str, idx: int, tone_cfg: Dict[str, Any], use_emojis: bool) -> str:
    clean = line.strip()
    if not clean:
        return ""

    is_bullet = clean.startswith(("- ", "* ", "• ", "▸ ", "– "))
    if is_bullet:
        clean = clean[2:].strip()

    if not is_bullet:
        return clean

    bullet = tone_cfg["bullet"]
    if use_emojis:
        return f"{_emoji_for_index(idx)} {clean}"
    return f"{bullet} {clean}"


def _style_body_paragraph(paragraph, theme: Dict[str, Any], tone_cfg: Dict[str, Any], is_bullet: bool) -> None:
    paragraph.font.name = theme["body_font"]
    paragraph.font.size = Pt(tone_cfg["body_size"])
    paragraph.font.color.rgb = _rgb(theme["body_color"])
    paragraph.level = 0 if not is_bullet else 1


def _chart_type(kind: str) -> XL_CHART_TYPE:
    mapping = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
    }
    return mapping.get(kind.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)


def _add_native_chart(slide, chart_config: Dict[str, Any], theme: Dict[str, Any]) -> None:
    categories = chart_config.get("categories", [])
    series = chart_config.get("series", [])

    if not categories or not series:
        return

    data = CategoryChartData()
    data.categories = categories

    for s in series:
        data.add_series(str(s.get("name", "Series")), list(s.get("values", [])))

    kind = str(chart_config.get("type", "column")).lower()
    chart = slide.shapes.add_chart(
        _chart_type(kind),
        Inches(float(chart_config.get("x", 6.4))),
        Inches(float(chart_config.get("y", 1.7))),
        Inches(float(chart_config.get("w", 6.2))),
        Inches(float(chart_config.get("h", 4.5))),
        data,
    ).chart

    title = chart_config.get("title")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(title)

    chart.has_legend = bool(chart_config.get("show_legend", True))
    try:
        for s in chart.series:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = _rgb(theme["accent"])
    except Exception:
        pass


def _add_matplotlib_chart(slide, chart_config: Dict[str, Any], theme: Dict[str, Any]) -> None:
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as e:
        raise RuntimeError(
            "matplotlib is required for chart mode 'matplotlib'. Install with: pip install matplotlib"
        ) from e

    categories = chart_config.get("categories", [])
    series = chart_config.get("series", [])
    if not categories or not series:
        return

    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
    fig.patch.set_facecolor("#" + theme["background"])
    ax.set_facecolor("#" + theme["background"])

    kind = str(chart_config.get("type", "column")).lower()

    if kind in {"bar", "column"}:
        width = 0.8 / max(1, len(series))
        for idx, s in enumerate(series):
            shift = (idx - (len(series) - 1) / 2) * width
            xs = [i + shift for i in range(len(categories))]
            ax.bar(xs, s.get("values", []), width=width, label=s.get("name", f"Series {idx + 1}"))
        ax.set_xticks(range(len(categories)))
        ax.set_xticklabels(categories)
    elif kind == "line":
        for idx, s in enumerate(series):
            ax.plot(categories, s.get("values", []), marker="o", label=s.get("name", f"Series {idx + 1}"))
    elif kind == "area":
        values = [list(s.get("values", [])) for s in series]
        labels = [str(s.get("name", f"Series {idx + 1}")) for idx, s in enumerate(series)]
        ax.stackplot(categories, *values, labels=labels, alpha=0.8)
    elif kind == "pie":
        first = series[0]
        ax.pie(first.get("values", []), labels=categories, autopct="%1.1f%%")
    else:
        ax.plot(categories, series[0].get("values", []), marker="o")

    ax.tick_params(colors="#" + theme["body_color"])
    for spine in ax.spines.values():
        spine.set_color("#" + theme["muted"])

    if chart_config.get("title"):
        ax.set_title(str(chart_config.get("title")), color="#" + theme["title_color"])

    if kind != "pie" and chart_config.get("show_legend", True):
        legend = ax.legend()
        if legend:
            for txt in legend.get_texts():
                txt.set_color("#" + theme["body_color"])

    fd, temp_png = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    fig.tight_layout()
    fig.savefig(temp_png, transparent=False)
    plt.close(fig)

    slide.shapes.add_picture(
        temp_png,
        Inches(float(chart_config.get("x", 6.4))),
        Inches(float(chart_config.get("y", 1.7))),
        Inches(float(chart_config.get("w", 6.2))),
        Inches(float(chart_config.get("h", 4.5))),
    )

    try:
        os.remove(temp_png)
    except OSError:
        pass


def _add_chart(slide, chart_config: Dict[str, Any], chart_mode: str, theme: Dict[str, Any]) -> None:
    mode = str(chart_config.get("mode", chart_mode)).lower()
    if mode == "auto":
        mode = "native"

    if mode == "native":
        _add_native_chart(slide, chart_config, theme)
    elif mode == "matplotlib":
        _add_matplotlib_chart(slide, chart_config, theme)
    else:
        raise ValueError(f"Unsupported chart mode: {mode}")


def create_pptx(title: str, slides: List[Dict[str, Any]], output_path: str, **kwargs) -> Dict[str, Any]:
    """Create a styled PowerPoint presentation with optional charts.

    New options:
    - theme: modern preset (midnight-luxe, aurora-glow, obsidian-slate, ivory-bloom, neon-velocity)
    - chart_mode: native | matplotlib | auto
    - use_emojis: bool
    - tone: classic-formal | boardroom | conversational | laid-back
    - require_preflight: bool (if True, requires all preference keys)
    """
    theme_name = str(kwargs.get("theme", "midnight-luxe"))
    chart_mode = str(kwargs.get("chart_mode", "native")).lower()
    use_emojis = bool(kwargs.get("use_emojis", False))
    tone = str(kwargs.get("tone", "boardroom")).lower()
    include_title_slide = bool(kwargs.get("include_title_slide", False))

    if kwargs.get("require_preflight", False):
        required = ["theme", "chart_mode", "use_emojis", "tone"]
        missing = [key for key in required if key not in kwargs]
        if missing:
            return {
                "success": False,
                "error": f"Missing preflight preferences: {', '.join(missing)}",
                "required_preferences": required,
            }

    theme = THEMES.get(theme_name, THEMES["midnight-luxe"])
    tone_cfg = TONE_CONFIG.get(tone, TONE_CONFIG["boardroom"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if title and include_title_slide:
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        _apply_background(title_slide, theme)
        _add_accent_bar(title_slide, theme)
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title
        _style_title(title_slide, "title", theme, tone_cfg)

    for slide_data in slides:
        layout_type = str(slide_data.get("layout", "content")).lower()
        has_chart = isinstance(slide_data.get("chart"), dict)

        if layout_type == "title":
            slide_layout = prs.slide_layouts[0]
        elif layout_type in {"title_only", "content_with_chart", "chart"} or has_chart:
            slide_layout = prs.slide_layouts[5]
        else:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)
        _apply_background(slide, theme)
        _add_accent_bar(slide, theme)

        if "title" in slide_data and slide.shapes.title:
            slide.shapes.title.text = str(slide_data["title"])

        _style_title(slide, layout_type, theme, tone_cfg)

        raw_content = slide_data.get("content", "")
        lines = raw_content if isinstance(raw_content, list) else str(raw_content).split("\n")

        if has_chart:
            text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.8))
            tf = text_box.text_frame
            tf.clear()
        elif "content" in slide_data and len(slide.placeholders) >= 2:
            tf = slide.placeholders[1].text_frame
            tf.clear()
        else:
            tf = None

        if tf is not None:
            para_index = 0
            for idx, line in enumerate(lines):
                normalized = _format_line(str(line), idx, tone_cfg, use_emojis)
                if not normalized:
                    continue

                paragraph = tf.paragraphs[0] if para_index == 0 else tf.add_paragraph()
                paragraph.text = normalized
                is_bullet = normalized.startswith(
                    (
                        tone_cfg["bullet"] + " ",
                        "🚀",
                        "📈",
                        "✅",
                        "💡",
                        "⚙️",
                        "🎯",
                        "📊",
                        "🧠",
                        "👉",
                    )
                )
                _style_body_paragraph(paragraph, theme, tone_cfg, is_bullet)
                para_index += 1

        if has_chart:
            _add_chart(slide, slide_data["chart"], chart_mode, theme)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return {
        "output_path": output_path,
        "file_size": os.path.getsize(output_path),
        "slides_count": len(prs.slides),
        "theme": theme_name,
        "chart_mode": chart_mode,
        "use_emojis": use_emojis,
        "tone": tone,
    }


def add_transition_effects(
    input_path: str,
    output_path: Optional[str] = None,
    effect: str = "fade",
    duration: float = 0.5,
    **kwargs,
) -> Dict[str, Any]:
    """Placeholder transition handler.

    python-pptx does not fully support transition effects in a stable API.
    This function keeps the call surface without silently breaking files.
    """
    _ = (effect, duration, kwargs)
    output_path = output_path or input_path

    prs = Presentation(input_path)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return {
        "output_path": output_path,
        "warning": "Transition effects are not implemented due to python-pptx API limitations.",
    }


def extract_text(input_path: str, **kwargs) -> Dict[int, str]:
    """Extract text from each slide.

    Returns a dict {slide_number: text}.
    """
    prs = Presentation(input_path)
    result: Dict[int, str] = {}

    for i, slide in enumerate(prs.slides, 1):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        result[i] = "\n".join(slide_text)

    return result
